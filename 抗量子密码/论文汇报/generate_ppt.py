"""
Generate a group-meeting presentation PPT for:
"Efficient Fully Homomorphic Encryption from (Standard) LWE"
Brakerski & Vaikuntanathan, FOCS 2011

Focus: Thorough presentation of theoretical logic and innovations.
"""

import copy
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree

# ── Paths ──────────────────────────────────────────────
TEMPLATE = r"C:\Users\14090\Documents\我的第一个知识库\抗量子密码\论文汇报\汇报模板.pptx"
OUTPUT   = r"C:\Users\14090\Documents\我的第一个知识库\抗量子密码\论文汇报\FHE_from_LWE_汇报.pptx"

# ── Color constants ─────────────────────────────────────
DARK_BLUE   = RGBColor(0x1B, 0x3A, 0x5C)
MEDIUM_BLUE = RGBColor(0x2C, 0x5F, 0x8A)
LIGHT_BLUE  = RGBColor(0x3A, 0x7C, 0xBF)
ACCENT_BLUE = RGBColor(0x4A, 0x90, 0xD9)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY   = RGBColor(0x33, 0x33, 0x33)
MED_GRAY    = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY  = RGBColor(0xE8, 0xEC, 0xF0)
RED_ACCENT  = RGBColor(0xC0, 0x39, 0x2B)
GREEN_ACCENT= RGBColor(0x27, 0xAE, 0x60)
ORANGE_ACC  = RGBColor(0xE6, 0x7E, 0x22)

# ── Font settings ───────────────────────────────────────
TITLE_FONT = "Microsoft YaHei"
BODY_FONT  = "Microsoft YaHei"
MATH_FONT  = "Consolas"
EN_FONT    = "Calibri"

def set_font(run, size_pt, bold=False, italic=False, color=BLACK, font_name=BODY_FONT):
    """Configure a run's font."""
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name

def add_textbox(slide, left, top, width, height, text_items, line_spacing=1.2):
    """
    Add a textbox with rich text. text_items is a list of dicts:
    {"text": str, "size": pt, "bold": bool, "italic": bool, "color": RGBColor, "font": str, "alignment": PP_ALIGN}
    """
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, para_items in enumerate(text_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if isinstance(para_items, dict):
            para_items = [para_items]

        for item in para_items:
            run = p.add_run()
            run.text = item.get("text", "")
            set_font(run,
                     item.get("size", 14),
                     item.get("bold", False),
                     item.get("italic", False),
                     item.get("color", BLACK),
                     item.get("font", BODY_FONT))

        p.alignment = para_items[0].get("alignment", PP_ALIGN.LEFT)
        if line_spacing:
            p.line_spacing = Pt(item.get("size", 14) * line_spacing)

    return txBox

def add_bullet_slide(slide, title_text, bullets, notes_text=None):
    """Add title + bullet content to a slide using layout 1."""
    # Set title
    if slide.shapes.title:
        slide.shapes.title.text = ""
        p = slide.shapes.title.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = title_text
        set_font(run, 32, bold=True, color=DARK_BLUE, font_name=TITLE_FONT)

    # Add bullet content in the body area
    body_ph = None
    for ph in slide.placeholders:
        if ph.placeholder_format.type == 7 or ph.placeholder_format.type == 2:  # OBJECT or BODY
            body_ph = ph
            break

    if body_ph:
        tf = body_ph.text_frame
        tf.clear()
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            if isinstance(bullet, dict):
                # Main bullet with sub-bullets
                run = p.add_run()
                run.text = bullet["main"]
                set_font(run, bullet.get("size", 16), bold=bullet.get("bold", False),
                        color=bullet.get("color", DARK_GRAY), font_name=BODY_FONT)
                p.level = bullet.get("level", 0)

                for sub in bullet.get("subs", []):
                    sp = tf.add_paragraph()
                    sr = sp.add_run()
                    sr.text = sub
                    set_font(sr, 13, bold=False, color=MED_GRAY, font_name=BODY_FONT)
                    sp.level = 1
            else:
                run = p.add_run()
                run.text = bullet
                set_font(run, 16, bold=False, color=DARK_GRAY, font_name=BODY_FONT)
                p.level = 0

            p.line_spacing = Pt(24)

    if notes_text and slide.has_notes_slide:
        slide.notes_slide.notes_text_frame.text = notes_text

def add_section_slide(slide, section_num, section_title_en, section_title_cn, notes=""):
    """Create a section divider slide."""
    # Use a blank layout and add decorative elements + text
    # Left colored bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(0.15), Inches(7.5)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = MEDIUM_BLUE
    bar.line.fill.background()

    # Section number
    add_textbox(slide, 1.0, 2.0, 3.0, 1.2, [
        [{"text": f"Part {section_num}", "size": 18, "bold": False, "color": MEDIUM_BLUE, "font": EN_FONT}]
    ])

    # English title
    add_textbox(slide, 1.0, 2.8, 11.0, 1.5, [
        [{"text": section_title_en, "size": 40, "bold": True, "color": DARK_BLUE, "font": EN_FONT}]
    ])

    # Chinese title
    add_textbox(slide, 1.0, 4.2, 11.0, 0.8, [
        [{"text": section_title_cn, "size": 22, "bold": False, "color": MED_GRAY, "font": TITLE_FONT}]
    ])

    # Thin separator line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1.0), Inches(5.2), Inches(4.0), Inches(0.03)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = MEDIUM_BLUE
    line.line.fill.background()

    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

def add_math_box(slide, left, top, width, height, math_text, bg_color=LIGHT_GRAY):
    """Add a highlighted math/formula box."""
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = bg_color
    box.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    box.line.width = Pt(0.5)

    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = math_text
    set_font(run, 13, bold=False, color=BLACK, font_name=MATH_FONT)
    p.alignment = PP_ALIGN.CENTER

    return box

def add_content_title(slide, title_text):
    """Add a formatted title to a blank slide."""
    add_textbox(slide, 0.9, 0.3, 11.5, 0.8, [
        [{"text": title_text, "size": 30, "bold": True, "color": DARK_BLUE, "font": TITLE_FONT}]
    ])
    # Underline
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.9), Inches(1.05), Inches(3.0), Inches(0.04)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = MEDIUM_BLUE
    line.line.fill.background()

def add_flow_arrow(slide, left, top, width, height):
    """Add a right arrow for flow diagrams."""
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = MEDIUM_BLUE
    arrow.line.fill.background()
    return arrow

def add_highlight_box(slide, left, top, width, height, text, color=ACCENT_BLUE):
    """Add a highlighted callout box."""
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()

    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    set_font(run, 14, bold=True, color=WHITE, font_name=BODY_FONT)
    p.alignment = PP_ALIGN.CENTER
    return box

def add_comparison_column(slide, left, top, width, height, title, items, color, is_good=True):
    """Add a comparison column (pros/cons)."""
    # Title
    add_textbox(slide, left, top, width, 0.5, [
        [{"text": title, "size": 18, "bold": True, "color": color, "font": TITLE_FONT,
          "alignment": PP_ALIGN.CENTER}]
    ])

    # Items
    y = top + 0.55
    for item in items:
        prefix = "✓ " if is_good else "✗ "
        add_textbox(slide, left + 0.1, y, width - 0.2, 0.45, [
            [{"text": f"{prefix}{item}", "size": 12, "bold": False, "color": DARK_GRAY, "font": BODY_FONT}]
        ])
        y += 0.42

# ═══════════════════════════════════════════════════════════
#  MAIN GENERATION
# ═══════════════════════════════════════════════════════════

prs = Presentation(TEMPLATE)

# ── Delete all existing slides (keep them in reverse order) ──
slide_ids = [slide.slide_id for slide in prs.slides]
for slide_id in slide_ids:
    # Delete by rId
    rId = None
    for r in prs.part.rel_rels:
        rel = prs.part.rel_rels[r]
        if hasattr(rel, 'target_part') and hasattr(rel.target_part, 'slide_id'):
            try:
                if rel.target_part.slide_id == slide_id:
                    rId = r
                    break
            except:
                pass

    # Alternative: use XML manipulation
    sldIdLst = prs.presentation.sldIdLst
    for sldId in list(sldIdLst):
        if sldId.get('id') == str(slide_id):
            sldIdLst.remove(sldId)
            break

# ── Helper: add a new slide from a layout ──
def new_slide(layout_idx):
    """Add a slide with given layout index."""
    layout = prs.slide_layouts[layout_idx]
    return prs.slides.add_slide(layout)

# ── Helper: delete all existing slides properly ──
def delete_all_slides(pres):
    """Delete all slides from presentation."""
    sldIdLst = pres.presentation.sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        if rId:
            try:
                pres.part.drop_rel(rId)
            except:
                pass
        sldIdLst.remove(sldId)

delete_all_slides(prs)

# ═══════════════════════════════════════════════════════════
#  SLIDE 1: TITLE SLIDE
# ═══════════════════════════════════════════════════════════
slide = new_slide(0)  # Title layout
if slide.shapes.title:
    slide.shapes.title.text = ""
    p = slide.shapes.title.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Efficient Fully Homomorphic Encryption\nfrom (Standard) LWE"
    set_font(run, 36, bold=True, color=DARK_BLUE, font_name=EN_FONT)
    p.alignment = PP_ALIGN.CENTER

# Authors
add_textbox(slide, 2.5, 5.2, 8.3, 0.6, [
    [{"text": "Zvika Brakerski (Weizmann)  ·  Vinod Vaikuntanathan (U. Toronto)",
      "size": 16, "bold": False, "color": DARK_GRAY, "font": EN_FONT, "alignment": PP_ALIGN.CENTER}]
])
add_textbox(slide, 2.5, 5.7, 8.3, 0.5, [
    [{"text": "FOCS 2011 · IEEE Symposium on Foundations of Computer Science",
      "size": 13, "bold": False, "color": MED_GRAY, "font": EN_FONT, "alignment": PP_ALIGN.CENTER}]
])
add_textbox(slide, 2.5, 6.3, 8.3, 0.5, [
    [{"text": "组会汇报 · 抗量子密码方向",
      "size": 14, "bold": False, "color": MEDIUM_BLUE, "font": TITLE_FONT, "alignment": PP_ALIGN.CENTER}]
])

# ═══════════════════════════════════════════════════════════
#  SLIDE 2: OUTLINE
# ═══════════════════════════════════════════════════════════
slide = new_slide(1)  # Title + Content
if slide.shapes.title:
    slide.shapes.title.text = ""
    p = slide.shapes.title.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "汇报大纲"
    set_font(run, 34, bold=True, color=DARK_BLUE, font_name=TITLE_FONT)

outline_items = [
    {"main": "1.  研究背景与动机 —— 全同态加密与Gentry蓝图「, "bold」: True, "size": 18, "color": DARK_BLUE},
    {"main": "2.  前人工作的局限 —— 理想格假设与Squashing范式「, "bold」: True, "size": 18, "color": DARK_BLUE},
    {"main": "3.  核心创新 I：Re-linearization（重线性化）「, "bold」: True, "size": 18, "color": MEDIUM_BLUE},
    {"subs": ["→ 从标准LWE构造Somewhat HE，摆脱理想格依赖"]},
    {"main": "4.  核心创新 II：Dimension-Modulus Reduction（维度-模缩减）「, "bold」: True, "size": 18, "color": MEDIUM_BLUE},
    {"subs": ["→ 替代Squashing，实现Bootstrapping，达到全同态"]},
    {"main": "5.  BTS方案完整构造与安全性分析「, "bold」: True, "size": 18, "color": DARK_BLUE},
    {"main": "6.  应用：近最优私有信息检索（PIR）「, "bold」: True, "size": 18, "color": DARK_BLUE},
    {"main": "7.  总结与讨论「, "bold」: True, "size": 18, "color": DARK_BLUE},
]
add_bullet_slide(slide, "", outline_items)

# ═══════════════════════════════════════════════════════════
#  SLIDE 3: SECTION — Background
# ═══════════════════════════════════════════════════════════
slide = new_slide(6)  # Blank layout
add_section_slide(slide, 1, "Background & Motivation", "研究背景与动机")

# ═══════════════════════════════════════════════════════════
#  SLIDE 4: What is FHE?
# ═══════════════════════════════════════════════════════════
slide = new_slide(6)  # Blank
add_content_title(slide, "什么是全同态加密（FHE）？")

# Definition box
add_math_box(slide, 0.9, 1.4, 11.5, 0.7,
    "Enc(m₁), Enc(m₂)  ──[ Eval(f) ]──►  Enc(f(m₁, m₂))    （对密文任意计算）")

# Key points
add_textbox(slide, 0.9, 2.4, 5.5, 3.5, [
    [{"text": "核心能力「, "size」: 20, "bold": True, "color": DARK_BLUE, "font": TITLE_FONT}],
    [{"text": ""}],
    [{"text": "• 在加密数据上执行任意复杂度的计算「, "size」: 15, "color": DARK_GRAY}],
    [{"text": "• 计算结果仍为加密形式，只有私钥持有者能解密「, "size」: 15, "color": DARK_GRAY}],
    [{"text": '• 从数学上实现「可委托计算」：数据加密 → 云端处理 → 返回加密结果', "size": 15, "color": DARK_GRAY}],
    [{"text": ""}],
    [{"text": "应用场景「, "size」: 20, "bold": True, "color": DARK_BLUE, "font": TITLE_FONT}],
    [{"text": ""}],
    [{"text": "• 云计算隐私保护 · 加密数据库查询「, "size」: 15, "color": DARK_GRAY}],
    [{"text": "• 私有信息检索 (PIR) · 安全多方计算「, "size」: 15, "color": DARK_GRAY}],
    [{"text": "• 机器学习推理隐私保护「, "size」: 15, "color": DARK_GRAY}],
])

# Timeline on right
add_textbox(slide, 7.2, 2.4, 5.5, 3.8, [
    [{"text": "发展历程「, "size」: 20, "bold": True, "color": DARK_BLUE, "font": TITLE_FONT}],
    [{"text": ""}],
    [{"text": "1978  Rivest, Adleman, Dertouzos 提出FHE概念「, "size」: 13, "color": DARK_GRAY, "font": EN_FONT}],
    [{"text": "       ↓  密码学「圣杯」，30年无解「, "size」: 12, "color": MED_GRAY}],
    [{"text": "2009  Gentry 首个候选方案 [STOC'09]", "size": 13, "bold": True, "color": RED_ACCENT, "font": EN_FONT}],
    [{"text": "       基于理想格 (Ideal Lattices)", "size": 12, "color": MED_GRAY}],
    [{"text": "2010  van Dijk et al. 整数上的FHE [Eurocrypt'10]", "size": 13, "color": DARK_GRAY, "font": EN_FONT}],
    [{"text": "2011  Brakerski & Vaikuntanathan 本文「, "size」: 13, "bold": True, "color": MEDIUM_BLUE, "font": EN_FONT}],
    [{"text": "       → 仅基于标准LWE假设的FHE", "size": 12, "color": MED_GRAY}],
    [{"text": "2011  Brakerski & V.  Ring-LWE FHE [Crypto'11]", "size": 13, "color": DARK_GRAY, "font": EN_FONT}],
])

# ═══════════════════════════════════════════════════════════
#  SLIDE 5: Gentry's Blueprint
# ═══════════════════════════════════════════════════════════
slide = new_slide(6)  # Blank
add_content_title(slide, "Gentry 的 Bootstrapping 蓝图（2009）")

# Three-step flow
# Step 1
add_highlight_box(slide, 0.9, 1.8, 3.5, 0.8, "Step 1: Somewhat HE", MEDIUM_BLUE)
add_textbox(slide, 0.9, 2.7, 3.5, 1.2, [
    [{"text": "构造一个可对密文进行有限次加法和乘法运算的「部分同态」加密方案「, "size」: 12, "color": DARK_GRAY}],
    [{"text": "能评估深度 ≤ D 的电路「, "size」: 12, "bold": True, "color": DARK_BLUE, "font": MATH_FONT}],
])

# Arrow
add_flow_arrow(slide, 4.6, 2.1, 0.6, 0.3)

# Step 2
add_highlight_box(slide, 5.4, 1.8, 3.5, 0.8, "Step 2: Bootstrappable?", MEDIUM_BLUE)
add_textbox(slide, 5.4, 2.7, 3.5, 1.2, [
    [{"text": "关键条件：解密电路 + 1次乘法 的深度 ≤ D", "size": 12, "color": DARK_GRAY}],
    [{"text": "即方案能「自举」地评估自身的解密过程「, "size」: 12, "color": DARK_GRAY}],
])

# Arrow
add_flow_arrow(slide, 9.1, 2.1, 0.6, 0.3)

# Step 3
add_highlight_box(slide, 9.9, 1.8, 3.0, 0.8, "Step 3: FHE!", GREEN_ACCENT)
add_textbox(slide, 9.9, 2.7, 3.0, 1.2, [
    [{"text": "Bootstrapping定理：部分同态 + 自举 = 全同态「, "size」: 12, "color": DARK_GRAY, "font": BODY_FONT}],
    [{"text": "可评估任意深度电路「, "size」: 12, "bold": True, "color": GREEN_ACCENT, "font": MATH_FONT}],
])

# Bottom: Diagram explanation
add_math_box(slide, 0.9, 4.3, 11.5, 2.5,
    "Bootstrapping 核心思想（简化版）：                                           \n"
    "                                                                             \n"
    "  给定: 密文 c = Enc(pk, m)，希望「刷新」噪声使之可以继续参与同态运算           \n"
    "                                                                             \n"
    "  Step A: 将 c 用新公钥 pk₂ 加密：c* = Enc(pk₂, c)   ← 双重加密              \n"
    "  Step B: 同态评估解密电路：Enc(pk₂, m) = Eval(Dec, c*, Enc(pk₂, sk₁))       \n"
    "           ↳ 方案用自己的同态能力「解开」内层加密，得到新公钥下的密文           \n"
    "                                                                             \n"
    "  结果: 得到同一明文在新密钥下的「新鲜」密文 → 可继续参与同态运算！              \n"
    "  ▸ 关键：解密电路的复杂度必须 ≤ 方案的同态能力                               ")

# ═══════════════════════════════════════════════════════════
#  SLIDE 6: Limitations of Prior Work
# ═══════════════════════════════════════════════════════════
slide = new_slide(6)  # Blank
add_content_title(slide, "前人工作的两大局限")

# Limitation 1
add_highlight_box(slide, 0.9, 1.6, 5.5, 0.65, "局限 1：依赖理想格（Ideal Lattices）假设", RED_ACCENT)
add_textbox(slide, 0.9, 2.4, 5.5, 2.0, [
    [{"text": "• Gentry方案的Somewhat HE基于理想格上的困难问题「, "size」: 13, "color": DARK_GRAY}],
    [{"text": "• 理想格 (Ideal Lattices) 自然地同时支持加法和乘法「, "size」: 13, "color": DARK_GRAY}],
    [{"text": "  原因：理想是环的子结构，对+和×封闭「, "size」: 12, "color": MED_GRAY, "font": EN_FONT}],
    [{"text": "  相比之下，一般格只对加法封闭「, "size」: 12, "color": MED_GRAY, "font": EN_FONT}],
    [{"text": "• 问题：理想格 = 「特殊品种」，我们对其了解甚少「, "size」: 13, "bold": True, "color": RED_ACCENT}],
    [{"text": "  一般格的研究更深入 (LLL, Ajtai, Micciancio...)", "size": 12, "color": MED_GRAY}],
])
# Arrow in the middle
add_textbox(slide, 6.6, 2.5, 0.5, 0.5, [
    [{"text": "VS", "size": 20, "bold": True, "color": RED_ACCENT, "font": EN_FONT, "alignment": PP_ALIGN.CENTER}]
])

# Our goal for limitation 1
add_textbox(slide, 7.2, 1.6, 5.5, 3.0, [
    [{"text": "本文目标 1：去掉理想格「, "size」: 18, "bold": True, "color": MEDIUM_BLUE, "font": TITLE_FONT}],
    [{"text": ""}],
    [{"text": "在标准 LWE 假设上构造 Somewhat HE", "size": 14, "bold": True, "color": DARK_BLUE}],
    [{"text": "• LWE 问题等价于一般格上最坏情况困难问题「, "size」: 13, "color": DARK_GRAY}],
    [{"text": "  (Regev'05, Peikert'09: worst-case → average-case)", "size": 11, "color": MED_GRAY, "font": EN_FONT}],
    [{"text": "• 新工具：Re-linearization（重线性化）「, "size」: 13, "bold": True, "color": MEDIUM_BLUE}],
])

# Limitation 2
add_highlight_box(slide, 0.9, 4.7, 5.5, 0.65, "局限 2：Squashing 范式 + 稀疏子集和假设", RED_ACCENT)
add_textbox(slide, 0.9, 5.45, 5.5, 2.0, [
    [{"text": "• 部分同态 → 全同态需要 Bootstrapping", "size": 13, "color": DARK_GRAY}],
    [{"text": "• 但解密电路复杂度 (degree ≈ max(n, log q)) 超过同态能力「, "size」: 13, "color": DARK_GRAY}],
    [{"text": "• Gentry 的 Squashing：人为降低解密复杂度「, "size」: 13, "color": DARK_GRAY}],
    [{"text": "  代价：引入「稀疏子集和" (Sparse Subset-Sum) 假设」, "size": 12, "bold": True, "color": RED_ACCENT}],
    [{"text": "• 这个额外假设是所有先前方案的「主要缺陷」「, "size」: 13, "color": RED_ACCENT}],
])

add_textbox(slide, 7.2, 4.7, 5.5, 3.0, [
    [{"text": "本文目标 2：去掉 Squashing", "size": 18, "bold": True, "color": MEDIUM_BLUE, "font": TITLE_FONT}],
    [{"text": ""}],
    [{"text": "不引入额外假设，直接实现 Bootstrapping", "size": 14, "bold": True, "color": DARK_BLUE}],
    [{"text": "• 新工具：Dimension-Modulus Reduction", "size": 13, "bold": True, "color": MEDIUM_BLUE}],
    [{"text": "  （维度-模缩减）「, "size」: 12, "color": MED_GRAY, "font": TITLE_FONT}],
    [{"text": "• 额外收获：密文非常短 → 高效的 PIR", "size": 13, "color": DARK_GRAY}],
])

# ═══════════════════════════════════════════════════════════
#  SLIDE 7: SECTION — Re-linearization
# ═══════════════════════════════════════════════════════════
slide = new_slide(6)  # Blank
add_section_slide(slide, 2, "Re-linearization", "核心创新 I：重线性化 —— 从 LWE 构造 Somewhat HE")

# ═══════════════════════════════════════════════════════════
#  SLIDE 8: LWE Encryption Basics
# ═══════════════════════════════════════════════════════════
slide = new_slide(6)  # Blank
add_content_title(slide, "起点：Regev 的 LWE 加密方案 [STOC'05]")

add_textbox(slide, 0.9, 1.4, 5.8, 2.5, [
    [{"text": "LWE (Learning With Errors) 假设「, "size」: 20, "bold": True, "color": DARK_BLUE, "font": TITLE_FONT}],
    [{"text": ""}],
    [{"text": "给定 (aᵢ, ⟨aᵢ,s⟩ + eᵢ)，区分其与均匀随机分布是困难的「, "size」: 13, "color": DARK_GRAY, "font": MATH_FONT}],
    [{"text": "  aᵢ ∈ Z_qⁿ (均匀随机),  s ∈ Z_qⁿ (秘密),  eᵢ ← χ (小噪声)", "size": 12, "color": MED_GRAY, "font": MATH_FONT}],
    [{"text": "• LWE 的困难性 ≈ 一般格上最坏情况的 Short Vector Problem", "size": 13, "color": DARK_GRAY}],
    [{"text": "  → 已知最优算法运行时间：exp(Ω̃(n))", "size": 12, "color": MED_GRAY, "font": MATH_FONT}],
])

add_textbox(slide, 7.2, 1.4, 5.5, 2.5, [
    [{"text": "Regev 加密方案（单比特）「, "size」: 20, "bold": True, "color": DARK_BLUE, "font": TITLE_FONT}],
    [{"text": ""}],
    [{"text": "Secret Key:  s ∈ Z_qⁿ", "size": 14, "color": DARK_GRAY, "font": MATH_FONT}],
    [{"text": "Enc(μ∈{0,1}):  选择 a ← Z_qⁿ, e ← χ", "size": 14, "color": DARK_GRAY, "font": MATH_FONT}],
    [{"text": "    c = (a,  b = ⟨a,s⟩ + 2e + μ)", "size": 14, "bold": True, "color": BLACK, "font": MATH_FONT}],
    [{"text": "Dec(c):  b - ⟨a,s⟩ mod q = 2e + μ", "size": 14, "color": DARK_GRAY, "font": MATH_FONT}],
    [{"text": "    再 mod 2 → μ  (e << q, 2e 不干扰)", "size": 14, "color": DARK_GRAY, "font": MATH_FONT}],
])

# Key observation box
add_math_box(slide, 0.9, 4.2, 11.5, 0.7,
    "核心观察: 两个「掩码」互不干扰 —— 秘密掩码 ⟨a,s⟩ 和 偶数掩码 2e 可以依次消除！")

add_textbox(slide, 0.9, 5.2, 11.5, 2.0, [
    [{"text": "解密视角的转换 —— 将密文看作关于密钥的线性函数「, "size」: 18, "bold": True, "color": DARK_BLUE, "font": TITLE_FONT}],
    [{"text": ""}],
    [{"text": "定义: 对密文 c = (a,b)，考虑函数  f_{a,b}(x) = b - ⟨a,x⟩  (mod q)", "size": 14, "color": DARK_GRAY, "font": MATH_FONT}],
    [{"text": "那么解密就是: 在 x = s 处求值 f，然后 mod 2 → μ", "size": 14, "color": DARK_GRAY, "font": MATH_FONT}],
    [{"text": "即:  Dec(s, c) = f_{a,b}(s) mod 2 = (b - ⟨a,s⟩ mod q) mod 2 = μ", "size": 14, "bold": True, "color": BLACK, "font": MATH_FONT}],
])

# ═══════════════════════════════════════════════════════════
#  SLIDE 9: The Multiplication Problem
# ═══════════════════════════════════════════════════════════
slide = new_slide(6)  # Blank
add_content_title(slide, "同态乘法难题：从线性到二次的爆炸")

# Left: Addition (easy)
add_highlight_box(slide, 0.9, 1.5, 5.5, 0.55, "同态加法：自然成立", GREEN_ACCENT)
add_textbox(slide, 0.9, 2.2, 5.5, 1.5, [
    [{"text": "f_{a,b}(x) + f_{a',b'}(x)", "size": 15, "color": BLACK, "font": MATH_FONT}],
    [{"text": "  = (b - ⟨a,x⟩) + (b' - ⟨a',x⟩)", "size": 14, "color": DARK_GRAY, "font": MATH_FONT}],
    [{"text": "  = (b+b') - ⟨a+a', x⟩", "size": 14, "color": DARK_GRAY, "font": MATH_FONT}],
    [{"text": "  = f_{a+a', b+b'}(x)    ← 仍是线性函数！「, "size」: 15, "bold": True, "color": GREEN_ACCENT, "font": MATH_FONT}],
    [{"text": "→ 密文 (a+a', b+b') 解密得到 μ₁+μ₂ ✓", "size": 13, "color": DARK_GRAY}],
])

# Right: Multiplication (problem)
add_highlight_box(slide, 7.2, 1.5, 5.5, 0.55, "同态乘法：从线性到二次 —— 密文膨胀！", RED_ACCENT)
add_textbox(slide, 7.2, 2.2, 5.5, 2.5, [
    [{"text": "f₁(x) · f₂(x) = (b - ⟨a,x⟩)(b' - ⟨a',x⟩)", "size": 14, "color": BLACK, "font": MATH_FONT}],
    [{"text": "展开后得到:", "size": 13, "color": DARK_GRAY}],
    [{"text": "  φ(x) = h₀ + Σᵢ hᵢ·x[i] + Σᵢⱼ hᵢⱼ·x[i]x[j]", "size": 14, "bold": True, "color": RED_ACCENT, "font": MATH_FONT}],
    [{"text": "  ↑ 二次项！不再是线性函数！「, "size」: 13, "color": RED_ACCENT}],
    [{"text": "", "size": 6}],
    [{"text": "密文大小:  n+1  →  ~n²/2  个系数  (爆炸!!)", "size": 15, "bold": True, "color": RED_ACCENT}],
    [{"text": "→ 一次乘法后密文就膨胀到不可用 [GHV10]", "size": 13, "color": DARK_GRAY}],
])

# Bottom: the core question
add_math_box(slide, 0.9, 5.1, 11.5, 1.6,
    "核心矛盾：                                                                  \n"
    "  同态加法 → 线性 + 线性 = 线性           ✓  自然支持                        \n"
    "  同态乘法 → 线性 × 线性 = 二次多项式     ✗  密文从 O(n) 膨胀为 O(n²)        \n"
    "                                                                             \n"
    "  ═══════════════  ═══════════════════════════════════════════════            \n"
    "  问题: 如何在做乘法时，把「二次」密文压缩回「线性」？                            \n"
    "  答案: Re-linearization（重线性化）—— 本论文的第一个核心创新！               ")

# ═══════════════════════════════════════════════════════════
#  SLIDE 10: Re-linearization Technical Details
# ═══════════════════════════════════════════════════════════
slide = new_slide(6)  # Blank
add_content_title(slide, "Re-linearization（重线性化）技术详解")

# Step-by-step
add_textbox(slide, 0.9, 1.4, 11.5, 0.8, [
    [{"text": "核心思想：发布「二次项在旧密钥下的密文" → 在新密钥下将其视为线性组合」, "size": 16, "bold": True, "color": DARK_BLUE}],
])

# Step 1
add_highlight_box(slide, 0.9, 2.3, 3.6, 0.5, "Step 1：密钥链", ACCENT_BLUE)
add_textbox(slide, 0.9, 2.9, 3.6, 1.5, [
    [{"text": "生成 L+1 层密钥：「, "size」: 13, "color": DARK_GRAY}],
    [{"text": "  s₀, s₁, s₂, ..., s_L", "size": 14, "bold": True, "color": BLACK, "font": MATH_FONT}],
    [{"text": "每个密钥 s_ℓ ∈ Z_qⁿ", "size": 12, "color": MED_GRAY, "font": MATH_FONT}],
    [{"text": "支持 L 层乘法运算「, "size」: 12, "color": MED_GRAY}],
])

# Step 2
add_highlight_box(slide, 4.8, 2.3, 3.6, 0.5, "Step 2：发布辅助参数", ACCENT_BLUE)
add_textbox(slide, 4.8, 2.9, 3.6, 1.5, [
    [{"text": "对于所有 i,j ∈ [n], τ ∈ [log q]:", "size": 12, "font": MATH_FONT}],
    [{"text": "b_{ℓ,i,j,τ} = ⟨a, s_ℓ⟩ + 2e", "size": 13, "font": MATH_FONT}],
    [{"text": "          + 2^τ · s_{ℓ-1}[i]·s_{ℓ-1}[j]", "size": 13, "bold": True, "color": BLACK, "font": MATH_FONT}],
    [{"text": "称为「伪加密」（pseudo-encryption）「, "size」: 12, "color": MED_GRAY}],
])

# Step 3
add_highlight_box(slide, 8.7, 2.3, 3.7, 0.5, "Step 3：重线性化运算", ACCENT_BLUE)
add_textbox(slide, 8.7, 2.9, 3.7, 1.5, [
    [{"text": "2^τ·s[i]s[j] ≈ b - ⟨a, s_ℓ⟩", "size": 13, "color": BLACK, "font": MATH_FONT}],
    [{"text": "因此 Σ h_{ij}·s[i]s[j]", "size": 13, "color": DARK_GRAY, "font": MATH_FONT}],
    [{"text": " ≈ Σ h_{ij,τ}·(b_{τ} - ⟨a_{τ}, s_ℓ⟩)", "size": 13, "bold": True, "color": MEDIUM_BLUE, "font": MATH_FONT}],
    [{"text": " = 线性函数 in s_ℓ ！「, "size」: 14, "bold": True, "color": GREEN_ACCENT, "font": MATH_FONT}],
])

# The magic
add_math_box(slide, 0.9, 4.7, 11.5, 2.2,
    "重线性化的「魔法」：                                                          \n"
    "                                                                             \n"
    "  输入:  两个密文 (a,b) 和 (a',b')，对应密钥 s_{ℓ-1}，明文 μ, μ'             \n"
    "         （它们代表的线性函数 f, f' 的乘积 是 二次函数 φ）                    \n"
    "                                                                             \n"
    "  输出:  一个新密文 (a*, b*)，对应密钥 s_ℓ，明文 μ·μ'                        \n"
    "         （它代表的线性函数 f* 在 s_ℓ 上等于 φ 在 s_{ℓ-1} 上的值）           \n"
    "                                                                             \n"
    "  关键代价: 每个乘法层需要发布 O(n² log q) 个辅助参数                         \n"
    "  ▸ 参数的二进制分解 (h_{ij,τ} ∈ {0,1}) 控制噪声增长                        ")

# ═══════════════════════════════════════════════════════════
#  SLIDE 11: Re-linearization Summary
# ═══════════════════════════════════════════════════════════
slide = new_slide(6)  # Blank
add_content_title(slide, "Re-linearization：成就与局限")

add_textbox(slide, 0.9, 1.5, 5.8, 4.5, [
    [{"text": "✅ 取得的突破「, "size」: 22, "bold": True, "color": GREEN_ACCENT, "font": TITLE_FONT}],
    [{"text": ""}],
    [{"text": "1. 从标准 LWE 构造出 Somewhat HE", "size": 16, "bold": True, "color": DARK_BLUE}],
    [{"text": "   → 摆脱了对理想格的依赖！「, "size」: 14, "color": MED_GRAY}],
    [{"text": "   → 安全性仅基于一般格的最坏情况困难性「, "size」: 14, "color": MED_GRAY}],
    [{"text": "", "size": 8}],
    [{"text": "2. 支持多项式深度的同态运算「, "size」: 16, "bold": True, "color": DARK_BLUE}],
    [{"text": "   • 乘法深度 L ≈ ε·log n (0<ε<1)", "size": 14, "color": DARK_GRAY, "font": MATH_FONT}],
    [{"text": "   • 等价于可评估度数 D ≈ n^ε 的多项式「, "size」: 14, "color": DARK_GRAY, "font": MATH_FONT}],
    [{"text": "   • 密文大小始终保持 O(n) 不变「, "size」: 14, "color": DARK_GRAY}],
    [{"text": "", "size": 8}],
    [{"text": "3. 技术通用性强「, "size」: 16, "bold": True, "color": DARK_BLUE}],
    [{"text": "   • 密钥切换 (Key Switching) 的雏形「, "size」: 14, "color": DARK_GRAY}],
    [{"text": "   • 后续工作 (BGV, GSW, CKKS...) 广泛使用「, "size」: 14, "color": DARK_GRAY}],
])

add_textbox(slide, 7.2, 1.5, 5.5, 4.5, [
    [{"text": "⚠️ 仍存在的局限「, "size」: 22, "bold": True, "color": ORANGE_ACC, "font": TITLE_FONT}],
    [{"text": ""}],
    [{"text": "这个方案还只是 Somewhat HE：「, "size」: 16, "bold": True, "color": DARK_BLUE}],
    [{"text": "", "size": 6}],
    [{"text": "• 同态运算深度受限于参数 n", "size": 15, "color": DARK_GRAY}],
    [{"text": "  → 无法评估任意深度的电路「, "size」: 14, "color": MED_GRAY}],
    [{"text": "", "size": 6}],
    [{"text": "• 要变成 FHE，必须能 Bootstrapping", "size": 15, "color": DARK_GRAY}],
    [{"text": "", "size": 6}],
    [{"text": "• 解密电路复杂度:", "size": 15, "color": DARK_GRAY}],
    [{"text": "  degree ≈ max(n, log q)", "size": 15, "bold": True, "color": RED_ACCENT, "font": MATH_FONT}],
    [{"text": "  但方案只能支持 degree ≈ n^ε", "size": 14, "color": MED_GRAY, "font": MATH_FONT}],
    [{"text": "  → n^ε << max(n, log q)  → 不能自举！「, "size」: 15, "bold": True, "color": RED_ACCENT}],
    [{"text": "", "size": 8}],
    [{"text": "需要新工具来缩小解密电路的复杂度...", "size": 16, "bold": True, "color": MEDIUM_BLUE}],
])

# ═══════════════════════════════════════════════════════════
#  SLIDE 12: SECTION — Dimension-Modulus Reduction
# ═══════════════════════════════════════════════════════════
slide = new_slide(6)  # Blank
add_section_slide(slide, 3, "Dimension-Modulus Reduction", "核心创新 II：维度-模缩减 —— 去掉 Squashing 实现 Bootstrapping")

# ═══════════════════════════════════════════════════════════
#  SLIDE 13: The Bootstrapping Bottleneck
# ═══════════════════════════════════════════════════════════
slide = new_slide(6)  # Blank
add_content_title(slide, "Bootstrapping 的瓶颈与解决思路")

add_textbox(slide, 0.9, 1.5, 5.8, 3.5, [
    [{"text": "瓶颈分析「, "size」: 20, "bold": True, "color": RED_ACCENT, "font": TITLE_FONT}],
    [{"text": ""}],
    [{"text": "Bootstrapping 需要方案能够同态评估自身的解密电路「, "size」: 14, "color": DARK_GRAY}],
    [{"text": ""}],
    [{"text": "解密操作: d = b - ⟨a,s⟩ (mod q), 然后 mod 2", "size": 14, "color": BLACK, "font": MATH_FONT}],
    [{"text": "写成关于 s 的多项式 (将 s 的每个比特视为变量):", "size": 13, "color": DARK_GRAY}],
    [{"text": "  degree ≥ max(n, log q)", "size": 15, "bold": True, "color": RED_ACCENT, "font": MATH_FONT}],
    [{"text": "  (n 来自内积, log q 来自模约简)", "size": 12, "color": MED_GRAY, "font": MATH_FONT}],
    [{"text": ""}],
    [{"text": "但 Somewhat HE 只能支持 degree ≈ n^ε (ε<1)", "size": 14, "color": DARK_GRAY}],
    [{"text": "→ max(n, log q) >> n^ε  无法满足自举条件！「, "size」: 14, "bold": True, "color": RED_ACCENT}],
])

add_textbox(slide, 7.2, 1.5, 5.5, 3.5, [
    [{"text": "两种解决路径「, "size」: 20, "bold": True, "color": DARK_BLUE, "font": TITLE_FONT}],
    [{"text": ""}],
    [{"text": "路径 A: Squashing (Gentry'09)", "size": 16, "bold": True, "color": MED_GRAY}],
    [{"text": "• 在公钥中加入关于私钥的「提示」「, "size」: 13, "color": DARK_GRAY}],
    [{"text": "• 人为降低解密多项式的次数「, "size」: 13, "color": DARK_GRAY}],
    [{"text": "• ❌ 引入稀疏子集和假设「, "size」: 14, "bold": True, "color": RED_ACCENT}],
    [{"text": "", "size": 8}],
    [{"text": "路径 B: Dimension-Modulus Reduction (本文)", "size": 16, "bold": True, "color": MEDIUM_BLUE}],
    [{"text": "• 将密文从大参数 (n,q) 转换到小参数 (k,p)", "size": 13, "color": DARK_GRAY}],
    [{"text": "• 小参数下解密多项式的 degree 显著降低「, "size」: 13, "color": DARK_GRAY}],
    [{"text": "• ✅ 不需要任何额外假设！「, "size」: 14, "bold": True, "color": GREEN_ACCENT}],
    [{"text": "• ✅ 密文变短，一举两得！「, "size」: 14, "bold": True, "color": GREEN_ACCENT}],
])

# Visual comparison
add_math_box(slide, 0.9, 5.3, 11.5, 1.8,
    "参数转换示意：                                                              \n"
    "                                                                             \n"
    "  原密文:   维度 n (大, 如 n = k⁴),  模数 q (大, 如 q ≈ 2^√n)               \n"
    "                                       ↓  Dimension-Modulus Reduction        \n"
    "  新密文:   维度 k (小, ≈ 安全参数),  模数 p (小, p = poly(k))               \n"
    "                                                                             \n"
    "  效果:  max(n, log q)  →  max(k, log p)   (后者小得多，可被同态评估！)      ")

# ═══════════════════════════════════════════════════════════
#  SLIDE 14: Dimension-Modulus Reduction Details
# ═══════════════════════════════════════════════════════════
slide = new_slide(6)  # Blank
add_content_title(slide, "维度-模缩减技术详解")

add_textbox(slide, 0.9, 1.4, 11.5, 0.7, [
    [{"text": "核心思路：利用密钥切换 (Key Switching) + 缩放近似，将密文参数 (n, q) 降到 (k, p)", "size": 15, "bold": True, "color": DARK_BLUE}],
])

# Step 1: Dimension reduction
add_highlight_box(slide, 0.9, 2.3, 5.5, 0.5, "Step 1: 维度缩减 (n → k)", ACCENT_BLUE)
add_textbox(slide, 0.9, 2.9, 5.5, 2.0, [
    [{"text": "回顾 Re-linearization 的思想:", "size": 13, "color": DARK_GRAY}],
    [{"text": "可以将对 s (n维) 的密文转换为对 t 的密文「, "size」: 13, "color": DARK_GRAY}],
    [{"text": "", "size": 5}],
    [{"text": "这次我们选择新的短密钥:", "size": 13, "color": DARK_GRAY}],
    [{"text": "  ŝ ∈ Z_p^k   (k << n)", "size": 15, "bold": True, "color": BLACK, "font": MATH_FONT}],
    [{"text": "", "size": 5}],
    [{"text": "关键: s 和 ŝ 的维度不需要相同！「, "size」: 14, "bold": True, "color": MEDIUM_BLUE}],
    [{"text": "  → Re-linearization 过程照样工作「, "size」: 13, "color": DARK_GRAY}],
    [{"text": "  → 从 (n, log q) 降到 (k, log q)", "size": 13, "color": DARK_GRAY, "font": MATH_FONT}],
])

# Step 2: Modulus reduction
add_highlight_box(slide, 7.2, 2.3, 5.5, 0.5, "Step 2: 模缩减 (q → p)", ACCENT_BLUE)
add_textbox(slide, 7.2, 2.9, 5.5, 2.5, [
    [{"text": "将 Z_q 中的元素缩放到 Z_p:", "size": 13, "color": DARK_GRAY}],
    [{"text": "", "size": 5}],
    [{"text": "scale(x) = ⌊(p/q) · x⌉  (四舍五入)", "size": 14, "bold": True, "color": BLACK, "font": MATH_FONT}],
    [{"text": "", "size": 5}],
    [{"text": "用于辅助参数的构造:", "size": 13, "color": DARK_GRAY}],
    [{"text": "  b̂_{i,τ} = ⟨â_{i,τ}, ŝ⟩ + ê", "size": 13, "color": BLACK, "font": MATH_FONT}],
    [{"text": "          + ⌊(p/q) · 2^τ · s_L[i]⌉", "size": 13, "bold": True, "color": MEDIUM_BLUE, "font": MATH_FONT}],
    [{"text": "", "size": 5}],
    [{"text": "缩放引入的舍入误差 ≤ 1/2，在噪声预算内可控「, "size」: 13, "color": DARK_GRAY}],
])

# Bottom summary
add_math_box(slide, 0.9, 5.4, 11.5, 1.5,
    "维度-模缩减的完整效果：                                                    \n"
    "                                                                             \n"
    "  输入密文:  c ∈ Z_q^n × Z_q   参数 (n, log q)    解密多项式 degree ≈ max(n, log q)\n"
    "       ↓  应用 Ψ̂ 中的辅助参数，进行一次"密钥切换 + 缩放"                    \n"
    "  输出密文:  ĉ ∈ Z_p^k × Z_p   参数 (k, log p)    解密多项式 degree ≈ max(k, log p)\n"
    "                                                                             \n"
    "  选 k = poly(κ), p = poly(k) → max(k, log p) 很小 → 解密可被同态评估！     ")

# ═══════════════════════════════════════════════════════════
#  SLIDE 15: BTS Scheme — Full Construction
# ═══════════════════════════════════════════════════════════
slide = new_slide(6)  # Blank
add_content_title(slide, "BTS 方案完整构造")

# KeyGen
add_textbox(slide, 0.9, 1.5, 5.8, 2.5, [
    [{"text": "BTS.KeyGen(1^κ) → (pk, evk, sk)", "size": 16, "bold": True, "color": DARK_BLUE, "font": EN_FONT}],
    [{"text": ""}],
    [{"text": "1. 采样「长」密钥链: s₀,...,s_L ← Z_q^n", "size": 13, "color": DARK_GRAY, "font": MATH_FONT}],
    [{"text": "2. 生成重线性化参数 Ψ:", "size": 13, "color": DARK_GRAY}],
    [{"text": "   对于 ℓ∈[L], i≤j≤n, τ∈[log q]:", "size": 12, "color": MED_GRAY, "font": MATH_FONT}],
    [{"text": "   b_{ℓ,i,j,τ} = ⟨a, s_ℓ⟩ + 2e + 2^τ·s_{ℓ-1}[i]·s_{ℓ-1}[j]", "size": 12, "color": BLACK, "font": MATH_FONT}],
    [{"text": "3. 采样「短」密钥: ŝ ← Z_p^k", "size": 13, "color": DARK_GRAY, "font": MATH_FONT}],
    [{"text": "4. 生成维度-模缩减桥接参数 Ψ̂:", "size": 13, "color": DARK_GRAY}],
    [{"text": "   对于 i∈[n], τ∈[log q]:", "size": 12, "color": MED_GRAY, "font": MATH_FONT}],
    [{"text": "   b̂_{i,τ} = ⟨â, ŝ⟩ + ê + ⌊(p/q)·2^τ·s_L[i]⌉", "size": 12, "color": BLACK, "font": MATH_FONT}],
    [{"text": "5. 生成公钥: A ← Z_q^{m×n}, b = A·s₀ + 2e", "size": 13, "color": DARK_GRAY, "font": MATH_FONT}],
    [{"text": "6. 输出: pk=(A,b), evk=(Ψ,Ψ̂), sk=ŝ", "size": 13, "bold": True, "color": DARK_BLUE, "font": MATH_FONT}],
])

# Enc/Dec/Eval
add_textbox(slide, 7.2, 1.5, 5.5, 2.5, [
    [{"text": "BTS.Enc(pk, μ):", "size": 16, "bold": True, "color": DARK_BLUE, "font": EN_FONT}],
    [{"text": "  选取 r ← {0,1}^m", "size": 13, "color": DARK_GRAY, "font": MATH_FONT}],
    [{"text": "  v = A^T·r,  w = b^T·r + μ", "size": 13, "color": BLACK, "font": MATH_FONT}],
    [{"text": "  输出: c = ((v,w), tag=0)", "size": 13, "bold": True, "color": DARK_BLUE, "font": MATH_FONT}],
    [{"text": ""}],
    [{"text": "BTS.Dec(ŝ, ĉ=(v̂,ŵ)):", "size": 16, "bold": True, "color": DARK_BLUE, "font": EN_FONT}],
    [{"text": "  μ* = (ŵ - ⟨v̂,ŝ⟩ mod p) mod 2", "size": 13, "color": BLACK, "font": MATH_FONT}],
    [{"text": ""}],
    [{"text": "BTS.Eval(evk, f, c₁,...,c_t):", "size": 16, "bold": True, "color": DARK_BLUE, "font": EN_FONT}],
    [{"text": "  Gate-by-gate 同态评估电路 f", "size": 13, "color": DARK_GRAY}],
    [{"text": "  最后一步: 维度-模缩减「, "size」: 13, "bold": True, "color": MEDIUM_BLUE}],
    [{"text": "  输出: ĉ ∈ Z_p^k × Z_p (短密文!)", "size": 13, "bold": True, "color": GREEN_ACCENT, "font": MATH_FONT}],
])

# Invariant
add_math_box(slide, 0.9, 5.3, 11.5, 1.6,
    "同态评估的核心不变量 (Invariant)：                                         \n"
    "  对于任意 tag=ℓ 的密文 c = ((v,w), ℓ)，始终满足:                          \n"
    "      w - ⟨v, s_ℓ⟩ ≡ μ + 2·e  (mod q)                                      \n"
    "  其中 μ 是明文，e 是小噪声。解密正确性依赖于 |e| < q/4。                   ")

# ═══════════════════════════════════════════════════════════
#  SLIDE 16: Homomorphic Evaluation Details
# ═══════════════════════════════════════════════════════════
slide = new_slide(6)  # Blank
add_content_title(slide, "同态加法与乘法的详细过程")

# Addition
add_textbox(slide, 0.9, 1.5, 5.5, 2.2, [
    [{"text": "同态加法 Gate", "size": 18, "bold": True, "color": GREEN_ACCENT, "font": TITLE_FONT}],
    [{"text": ""}],
    [{"text": "输入: c₁=((v₁,w₁),ℓ), ..., c_t=((v_t,w_t),ℓ)", "size": 12, "color": DARK_GRAY, "font": MATH_FONT}],
    [{"text": "输出: c_add = ((Σvᵢ, Σwᵢ), ℓ)", "size": 13, "bold": True, "color": BLACK, "font": MATH_FONT}],
    [{"text": ""}],
    [{"text": "验证不变量:", "size": 13, "color": DARK_BLUE}],
    [{"text": "  (Σwᵢ) - ⟨Σvᵢ, s_ℓ⟩ = Σ(wᵢ - ⟨vᵢ, s_ℓ⟩)", "size": 12, "color": DARK_GRAY, "font": MATH_FONT}],
    [{"text": "  = Σ(μᵢ + 2eᵢ) = (Σμᵢ) + 2(Σeᵢ) ✓", "size": 12, "color": BLACK, "font": MATH_FONT}],
    [{"text": "→ 噪声线性增长，完全可控「, "size」: 13, "color": GREEN_ACCENT}],
])

# Multiplication
add_textbox(slide, 7.0, 1.5, 5.8, 4.0, [
    [{"text": "同态乘法 Gate", "size": 18, "bold": True, "color": RED_ACCENT, "font": TITLE_FONT}],
    [{"text": ""}],
    [{"text": "输入: c=((v,w),ℓ), c'=((v',w'),ℓ)  [同层级]", "size": 11, "color": DARK_GRAY, "font": MATH_FONT}],
    [{"text": ""}],
    [{"text": "Step 1: 构造二次型 Φ(x) = (w-⟨v,x⟩)(w'-⟨v',x⟩)", "size": 12, "color": DARK_BLUE, "font": MATH_FONT}],
    [{"text": "  = Σ_{i≤j≤n} h_{i,j} · x[i]x[j]", "size": 12, "color": BLACK, "font": MATH_FONT}],
    [{"text": ""}],
    [{"text": "Step 2: 比特分解 h_{i,j} = Σ_τ h_{i,j,τ}·2^τ", "size": 12, "color": DARK_BLUE, "font": MATH_FONT}],
    [{"text": "  → h_{i,j,τ} ∈ {0,1}，控制噪声增长「, "size」: 11, "color": MED_GRAY}],
    [{"text": ""}],
    [{"text": "Step 3: 用 Ψ 中参数重线性化:", "size": 12, "color": DARK_BLUE}],
    [{"text": "  v_mult = Σ h_{i,j,τ} · a_{ℓ+1,i,j,τ}", "size": 12, "color": BLACK, "font": MATH_FONT}],
    [{"text": "  w_mult = Σ h_{i,j,τ} · b_{ℓ+1,i,j,τ}", "size": 12, "color": BLACK, "font": MATH_FONT}],
    [{"text": ""}],
    [{"text": "输出: c_mult = ((v_mult, w_mult), ℓ+1)", "size": 12, "bold": True, "color": MEDIUM_BLUE, "font": MATH_FONT}],
    [{"text": "→ 解密得到 μ·μ' + 2e_mult  ✓", "size": 12, "color": GREEN_ACCENT}],
])

# Final D-M Reduction
add_math_box(slide, 0.9, 5.3, 11.5, 1.6,
    "最后一步 —— 维度-模缩减 (在电路顶层 ℓ=L 执行)：                             \n"
    "  φ(x) = ⌊(p/q)·((q+1)/2)·(w - ⟨v,x⟩)⌉  mod p                             \n"
    "  v̂ = 2·Σ h_{i,τ}·â_{i,τ}  mod p     ŵ = 2·Σ h_{i,τ}·b̂_{i,τ}  mod p     \n"
    "  输出 ĉ = (v̂, ŵ) ∈ Z_p^k × Z_p  —— k+1 个元素，每元素 log p 比特           ")

# ═══════════════════════════════════════════════════════════
#  SLIDE 17: Bootstrapping Achieved & Security
# ═══════════════════════════════════════════════════════════
slide = new_slide(6)  # Blank
add_content_title(slide, "Bootstrapping 达成与安全性")

add_textbox(slide, 0.9, 1.5, 5.8, 3.5, [
    [{"text": "参数设置（以运行示例说明）「, "size」: 18, "bold": True, "color": DARK_BLUE, "font": TITLE_FONT}],
    [{"text": ""}],
    [{"text": "  k = κ  (安全参数)", "size": 14, "color": DARK_GRAY, "font": MATH_FONT}],
    [{"text": "  n = k⁴  (长维度)", "size": 14, "color": DARK_GRAY, "font": MATH_FONT}],
    [{"text": "  q ≈ 2^{√n}  (大模数)", "size": 14, "color": DARK_GRAY, "font": MATH_FONT}],
    [{"text": "  p = poly(k)  (小模数)", "size": 14, "color": DARK_GRAY, "font": MATH_FONT}],
    [{"text": "  L = (1/3)log n = (4/3)log k  (乘法层数)", "size": 14, "color": DARK_GRAY, "font": MATH_FONT}],
    [{"text": ""}],
    [{"text": "Bootstrapping 条件验证:", "size": 16, "bold": True, "color": MEDIUM_BLUE}],
    [{"text": "  • BTS 可评估深度 ≈ n^ε = k^{c·ε} 的电路「, "size」: 13, "color": DARK_GRAY}],
    [{"text": "  • 解密电路深度 = O(log k + log log p)", "size": 13, "color": DARK_GRAY, "font": MATH_FONT}],
    [{"text": "  • 选取 c 足够大 → n^ε > 解密电路深度 ✓", "size": 14, "bold": True, "color": GREEN_ACCENT}],
])

add_textbox(slide, 7.2, 1.5, 5.5, 3.5, [
    [{"text": "安全性分析「, "size」: 18, "bold": True, "color": DARK_BLUE, "font": TITLE_FONT}],
    [{"text": ""}],
    [{"text": "方案安全性基于两个 LWE 假设:", "size": 14, "color": DARK_GRAY}],
    [{"text": ""}],
    [{"text": "1. DLWE_{n,q,χ}:  长参数下的 LWE", "size": 14, "color": DARK_BLUE, "font": MATH_FONT}],
    [{"text": "   用于公钥和重线性化参数的安全性「, "size」: 12, "color": MED_GRAY}],
    [{"text": ""}],
    [{"text": "2. DLWE_{k,p,χ̂}:  短参数下的 LWE", "size": 14, "color": DARK_BLUE, "font": MATH_FONT}],
    [{"text": "   用于维度-模缩减桥接参数的安全性「, "size」: 12, "color": MED_GRAY}],
    [{"text": ""}],
    [{"text": "虽然 k,p << n,q，但短参数方案不需要支持「, "size」: 13, "color": DARK_GRAY}],
    [{"text": "同态运算，因此可以使用更大的相对噪声，「, "size」: 13, "color": DARK_GRAY}],
    [{"text": "使得两个假设的困难度相当。「, "size」: 13, "color": DARK_GRAY}],
    [{"text": ""}],
    [{"text": "→ 仅依赖 LWE 假设 + 弱循环安全性 (Circular Security)", "size": 14, "bold": True, "color": MEDIUM_BLUE}],
])

# Theorem box
add_math_box(slide, 0.9, 5.4, 11.5, 1.2,
    "定理 (Brakerski-Vaikuntanathan, FOCS 2011):                                    \n"
    "  基于 DLWE_{n,q,χ} + DLWE_{k,p,χ̂} 假设，BTS 是一个可自举的 (Bootstrappable) 加密方案。\n"
    "  应用 Gentry 的自举定理 → 得到 Leveled FHE。  (若额外假设弱循环安全 → 得到纯 FHE)")

# ═══════════════════════════════════════════════════════════
#  SLIDE 18: SECTION — PIR
# ═══════════════════════════════════════════════════════════
slide = new_slide(6)  # Blank
add_section_slide(slide, 4, "Application: Near-Optimal PIR", "应用：近最优私有信息检索")

# ═══════════════════════════════════════════════════════════
#  SLIDE 19: PIR Protocol
# ═══════════════════════════════════════════════════════════
slide = new_slide(6)  # Blank
add_content_title(slide, "从短密文到高效私有信息检索 (PIR)")

add_textbox(slide, 0.9, 1.5, 5.8, 3.5, [
    [{"text": "PIR 问题设定「, "size」: 18, "bold": True, "color": DARK_BLUE, "font": TITLE_FONT}],
    [{"text": ""}],
    [{"text": "• 发送方持有大型数据库 DB (大小 N)", "size": 14, "color": DARK_GRAY}],
    [{"text": "• 接收方想查询第 i 条记录，但不泄露 i", "size": 14, "color": DARK_GRAY}],
    [{"text": "• 目标: 通信复杂度远小于 O(N)", "size": 14, "bold": True, "color": DARK_BLUE}],
    [{"text": "  最优目标: polylog(N)", "size": 13, "color": MED_GRAY, "font": MATH_FONT}],
    [{"text": ""}],
    [{"text": "FHE → PIR 的朴素方法:", "size": 15, "bold": True, "color": DARK_BLUE}],
    [{"text": "  1. 接收方加密索引: Enc(i)", "size": 13, "color": DARK_GRAY}],
    [{"text": "  2. 发送方同态评估 DB 访问函数「, "size」: 13, "color": DARK_GRAY}],
    [{"text": "  3. 发送方返回 Enc(DB[i])", "size": 13, "color": DARK_GRAY}],
    [{"text": "  问题: 需要逐比特加密索引 → O(log² N) 通信「, "size」: 12, "color": RED_ACCENT}],
])

add_textbox(slide, 7.2, 1.5, 5.5, 3.5, [
    [{"text": "本文的改进策略「, "size」: 18, "bold": True, "color": MEDIUM_BLUE, "font": TITLE_FONT}],
    [{"text": ""}],
    [{"text": "1. 用对称加密 (短密文) 加密查询索引「, "size」: 14, "color": DARK_GRAY}],
    [{"text": "   → 查询复杂度: O(log N) 比特「, "size」: 13, "color": GREEN_ACCENT}],
    [{"text": ""}],
    [{"text": "2. 将对称密钥用 FHE 加密放在公钥中「, "size」: 14, "color": DARK_GRAY}],
    [{"text": "   → 发送方可自行转换为同态密文「, "size」: 13, "color": MED_GRAY}],
    [{"text": ""}],
    [{"text": "3. BTS 的密文非常短:", "size": 14, "color": DARK_GRAY}],
    [{"text": "   |ĉ| = (k+1) log p = O(k log k) 比特「, "size」: 14, "bold": True, "color": BLACK, "font": MATH_FONT}],
    [{"text": "   → 响应复杂度: log N · polyloglog N", "size": 14, "bold": True, "color": GREEN_ACCENT}],
    [{"text": ""}],
    [{"text": "结果: 首个基于 LWE 的多对数级 PIR！「, "size」: 15, "bold": True, "color": MEDIUM_BLUE}],
    [{"text": "通信复杂度接近理论下界 log N", "size": 13, "color": DARK_GRAY}],
])

# Communication comparison
add_math_box(slide, 0.9, 5.3, 11.5, 1.5,
    "通信复杂度对比 (单比特查询, 公钥模型):                                      \n"
    "  朴素 FHE-PIR:     O(log² N)                                                \n"
    "  Gentry's SWHE:    O(log³ N)                                                \n"
    "  本文方案:          O(log N · polyloglog N)  ← 近最优！                      \n"
    "  理论下界:          Ω(log N)                                                 ")

# ═══════════════════════════════════════════════════════════
#  SLIDE 20: Summary & Contributions
# ═══════════════════════════════════════════════════════════
slide = new_slide(6)  # Blank
add_content_title(slide, "总结：主要贡献与影响")

# Contribution 1
add_highlight_box(slide, 0.9, 1.6, 5.5, 0.65, "贡献 1：Re-linearization 技术", MEDIUM_BLUE)
add_textbox(slide, 0.9, 2.4, 5.5, 2.2, [
    [{"text": "✓ 在标准 LWE 假设上构造 Somewhat HE", "size": 14, "bold": True, "color": GREEN_ACCENT}],
    [{"text": "  → 摆脱了对理想格 (Ideal Lattices) 的依赖「, "size」: 13, "color": DARK_GRAY}],
    [{"text": "✓ 新范式: 密钥切换 (Key Switching)", "size": 14, "bold": True, "color": GREEN_ACCENT}],
    [{"text": "  → 将二次密文压回线性，控制尺寸增长「, "size」: 13, "color": DARK_GRAY}],
    [{"text": "✓ 成为后续所有主流FHE方案的核心组件:", "size": 14, "bold": True, "color": GREEN_ACCENT}],
    [{"text": "  → BGV [Brakerski-Gentry-V. 2014]", "size": 12, "color": MED_GRAY, "font": EN_FONT}],
    [{"text": "  → GSW [Gentry-Sahai-Waters 2013]", "size": 12, "color": MED_GRAY, "font": EN_FONT}],
    [{"text": "  → CKKS [Cheon-Kim-Kim-Song 2017]", "size": 12, "color": MED_GRAY, "font": EN_FONT}],
])

# Contribution 2
add_highlight_box(slide, 7.2, 1.6, 5.5, 0.65, "贡献 2：Dimension-Modulus Reduction", MEDIUM_BLUE)
add_textbox(slide, 7.2, 2.4, 5.5, 2.2, [
    [{"text": "✓ 替代 Gentry 的 Squashing 范式「, "size」: 14, "bold": True, "color": GREEN_ACCENT}],
    [{"text": "  → 去掉稀疏子集和 (Sparse Subset-Sum) 假设「, "size」: 13, "color": DARK_GRAY}],
    [{"text": "✓ 实现纯 LWE 假设上的 Bootstrapping", "size": 14, "bold": True, "color": GREEN_ACCENT}],
    [{"text": "  → 通过缩减解密电路参数来满足自举条件「, "size」: 13, "color": DARK_GRAY}],
    [{"text": "✓ 额外收获: 非常短的密文「, "size」: 14, "bold": True, "color": GREEN_ACCENT}],
    [{"text": "  → 直接推动近最优 PIR 协议的构造「, "size」: 13, "color": DARK_GRAY}],
    [{"text": "✓ 维度-模缩减 → 模数切换 (Modulus Switching)", "size": 14, "bold": True, "color": GREEN_ACCENT}],
    [{"text": "  → 同样被后续工作广泛采用「, "size」: 13, "color": DARK_GRAY}],
])

# Open problems
add_textbox(slide, 0.9, 5.0, 11.5, 1.8, [
    [{"text": "开放问题与后续方向「, "size」: 18, "bold": True, "color": DARK_BLUE, "font": TITLE_FONT}],
    [{"text": ""}],
    [{"text": "• 循环安全性 (Circular Security): 要获得纯FHE (非Leveled) 需要额外的循环安全假设，能否去掉？「, "size」: 13, "color": DARK_GRAY}],
    [{"text": "• 常数摊销通信的 PIR: 本文实现了 logN·polyloglogN，能否降到 O(log N) 甚至常数摊销？「, "size」: 13, "color": DARK_GRAY}],
    [{"text": "• 实际效率: 参数仍然很大 (n=k⁴)，如何进一步缩小以实用化？ → 后续 BGV、GSW、CKKS 等工作逐步推进「, "size」: 13, "color": DARK_GRAY}],
])

# ═══════════════════════════════════════════════════════════
#  SLIDE 21: Thank You
# ═══════════════════════════════════════════════════════════
slide = new_slide(0)  # Title layout (for centered text)
if slide.shapes.title:
    slide.shapes.title.text = ""
    p = slide.shapes.title.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Thank You!  ·  感谢聆听"
    set_font(run, 40, bold=True, color=DARK_BLUE, font_name=EN_FONT)
    p.alignment = PP_ALIGN.CENTER

add_textbox(slide, 2.5, 4.8, 8.3, 1.5, [
    [{"text": "Efficient Fully Homomorphic Encryption from (Standard) LWE",
      "size": 16, "bold": False, "color": MED_GRAY, "font": EN_FONT, "alignment": PP_ALIGN.CENTER}],
    [{"text": "Brakerski & Vaikuntanathan · FOCS 2011",
      "size": 14, "bold": False, "color": MED_GRAY, "font": EN_FONT, "alignment": PP_ALIGN.CENTER}],
    [{"text": "",
      "size": 8}],
    [{"text": "Questions & Discussion",
      "size": 22, "bold": True, "color": MEDIUM_BLUE, "font": EN_FONT, "alignment": PP_ALIGN.CENTER}],
])

# ═══════════════════════════════════════════════════════════
#  SAVE
# ═══════════════════════════════════════════════════════════
prs.save(OUTPUT)
print(f"DONE! Saved to: {OUTPUT}")
print(f"Total slides: {len(prs.slides)}")
