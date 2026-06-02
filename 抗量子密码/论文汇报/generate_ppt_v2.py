# -*- coding: utf-8 -*-
"""
Generate FHE PPT from template by PRESERVING all visual elements
and ONLY modifying text content. Clones slides as needed.
"""
import copy, shutil, os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree

TEMPLATE = r"C:\Users\14090\Documents\我的第一个知识库\抗量子密码\论文汇报\汇报模板.pptx"
OUTPUT   = r"C:\Users\14090\Documents\我的第一个知识库\抗量子密码\论文汇报\FHE_from_LWE_汇报.pptx"
shutil.copy2(TEMPLATE, OUTPUT)
prs = Presentation(OUTPUT)

NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'
RNS = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'

DARK_BLUE  = RGBColor(0x1B, 0x3A, 0x5C)
MED_BLUE   = RGBColor(0x2C, 0x5F, 0x8A)
BLACK      = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY  = RGBColor(0x33, 0x33, 0x33)
MED_GRAY   = RGBColor(0x66, 0x66, 0x66)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GREEN      = RGBColor(0x27, 0xAE, 0x60)
RED        = RGBColor(0xC0, 0x39, 0x2B)

TFONT = "Microsoft YaHei"
MFONT = "Consolas"
EFONT = "Calibri"

# ================================================================
# HELPERS
# ================================================================

def find_shape(slide, name):
    for s in slide.shapes:
        if s.name == name:
            return s
    return None

def set_text(shape, text, size=14, bold=False, color=BLACK, font=TFONT, align=None):
    """Replace all text in a shape."""
    if not shape or not shape.has_text_frame:
        return
    tf = shape.text_frame
    for p in tf.paragraphs:
        p.clear()
    p = tf.paragraphs[0]
    run = p.add_run(); run.text = text
    run.font.size = Pt(size); run.font.bold = bold
    if color: run.font.color.rgb = color
    run.font.name = font
    if align: p.alignment = align

def set_multiline(shape, lines):
    """Set multiple lines of text in a shape. lines: list of (text, size, bold, color, font)."""
    if not shape or not shape.has_text_frame:
        return
    tf = shape.text_frame
    # Clear existing paragraphs
    for p in tf.paragraphs:
        p.clear()
    # Remove extra paragraphs
    ns = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
    for i, (text, size, bold, color, font) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run(); run.text = text
        run.font.size = Pt(size); run.font.bold = bold
        if color: run.font.color.rgb = color
        run.font.name = font

def clear_shape(shape):
    """Clear all text from a shape."""
    if shape and shape.has_text_frame:
        for p in shape.text_frame.paragraphs:
            p.clear()

def add_tb(slide, left, top, width, height, text, size=14, bold=False, color=BLACK, font=TFONT):
    """Add a simple text box."""
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; run = p.add_run(); run.text = text
    run.font.size = Pt(size); run.font.bold = bold
    run.font.color.rgb = color; run.font.name = font
    return tb

def add_multiline_tb(slide, left, top, width, height, lines):
    """Add a multiline textbox. lines: list of (text, size, bold, color, font)."""
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame; tf.word_wrap = True
    for i, (text, size, bold, color, font) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run(); run.text = text
        run.font.size = Pt(size); run.font.bold = bold
        run.font.color.rgb = color; run.font.name = font
        p.line_spacing = Pt(size * 1.3)
    return tb

def clone_slide(prs, src_idx):
    """Clone slide at src_idx and append to end. Return new slide index."""
    src = prs.slides[src_idx]
    layout = src.slide_layout
    # Find the slide's rId
    src_rId = None
    for rId, rel in src.part.rel_rels.items():
        if hasattr(rel, 'target_part'):
            pass

    # Use internal method: duplicate slide XML
    slide_xml = copy.deepcopy(src._element)
    # Get sldIdLst
    pres_elem = prs.part._element
    sldIdLst = pres_elem.find(f'.//{{{NS}}}sldIdLst')

    # Create new slide
    new_slide = prs.slides.add_slide(layout)

    # Copy shapes from source to new (remove defaults first)
    for shape in list(new_slide.shapes):
        shape._element.getparent().remove(shape._element)
    for shape in src.shapes:
        el = copy.deepcopy(shape._element)
        new_slide.shapes._spTree.append(el)

    return len(prs.slides) - 1

def section_title(slide, en_title, cn_title=""):
    """Update a section divider slide's title text."""
    # The section divider has TextBox 4 with the big English title
    tb4 = find_shape(slide, "TextBox 4")
    if tb4:
        set_text(tb4, en_title, size=40, bold=True, color=DARK_BLUE, font=EFONT)
    # The footer text
    tb6 = find_shape(slide, "文本框 6")
    if tb6:
        set_text(tb6, cn_title if cn_title else en_title, size=12, color=MED_GRAY, font=TFONT)

def content_title(slide, title_text):
    """Update a content slide's title bar (组合 3)."""
    g3 = find_shape(slide, "组合 3")
    if g3 and g3.shape_type == 6:  # Group
        for child in g3.shapes:
            if child.has_text_frame:
                set_text(child, title_text, size=26, bold=True, color=DARK_BLUE, font=TFONT)
                break

def clear_content_area(slide):
    """Clear all text from content area shapes (but keep sidebar table and title bar)."""
    keep_names = {"table 482", "组合 3", "Picture 4", "图片 10"}
    for shape in slide.shapes:
        if shape.name not in keep_names and shape.has_text_frame:
            clear_shape(shape)

# ================================================================
# MAP TEMPLATE SLIDES TO FHE CONTENT
# ================================================================

# SLIDE 1: TITLE - Update to FHE paper
slide = prs.slides[0]
set_text(find_shape(slide, "文本框 2"),
    "Efficient Fully Homomorphic Encryption\nfrom (Standard) LWE",
    size=28, bold=True, color=DARK_BLUE, font=EFONT, align=PP_ALIGN.CENTER)
set_text(find_shape(slide, "文本框 4"),
    "Zvika Brakerski (Weizmann)  |  Vinod Vaikuntanathan (U. Toronto)\nFOCS 2011  |  抗量子密码组会汇报",
    size=13, color=DARK_GRAY, font=TFONT, align=PP_ALIGN.CENTER)

# SLIDE 2: Section - "Introduction" -> "Background"
slide = prs.slides[1]
section_title(slide, "Background & Motivation", "Part 1 - 研究背景与动机")

# SLIDE 3: Content - What is FHE?
slide = prs.slides[2]
content_title(slide, "什么是全同态加密 (FHE) ?")
clear_content_area(slide)
add_multiline_tb(slide, 2.8, 1.6, 9.0, 5.0, [
    ("核心定义", 20, True, DARK_BLUE, TFONT),
    ("", 8, False, BLACK, TFONT),
    ("Enc(m1), Enc(m2)  --[ Eval(f) ]-->  Enc(f(m1,m2))", 15, True, BLACK, MFONT),
    ("在加密数据上执行任意计算，结果仍是密文", 14, False, DARK_GRAY, TFONT),
    ("", 8, False, BLACK, TFONT),
    ("密码学圣杯：1978年 Rivest-Adleman-Dertouzos 提出 → 2009年 Gentry 首个方案", 13, False, DARK_GRAY, TFONT),
    ("应用：云计算隐私 / 加密数据库查询 / 私有信息检索(PIR) / 安全多方计算", 13, False, DARK_GRAY, TFONT),
    ("", 8, False, BLACK, TFONT),
    ("Gentry方案：基于理想格(Ideal Lattices) → 需要Squashing+稀疏子集和假设", 13, False, MED_GRAY, TFONT),
    ("本文目标：仅基于标准LWE假设构造FHE，去掉理想格和稀疏子集和依赖", 14, True, MED_BLUE, TFONT),
])

# SLIDE 4: Content - Gentry's Blueprint
slide = prs.slides[3]
content_title(slide, "Gentry 的 Bootstrapping 蓝图")
clear_content_area(slide)
add_multiline_tb(slide, 2.8, 1.6, 9.0, 5.0, [
    ("三步构建全同态加密：", 18, True, DARK_BLUE, TFONT),
    ("", 8, False, BLACK, TFONT),
    ("Step 1: 构造 Somewhat Homomorphic Encryption (SWHE)", 15, True, MED_BLUE, TFONT),
    ("  可评估有限深度D的电路，支持加密数据的加法和乘法", 13, False, DARK_GRAY, TFONT),
    ("", 6, False, BLACK, TFONT),
    ("Step 2: 检查 Bootstrappable 条件", 15, True, MED_BLUE, TFONT),
    ("  解密电路的深度 + 1次额外乘法 <= D ?", 13, False, DARK_GRAY, TFONT),
    ("  即：方案能否同态评估自身的解密过程？", 13, False, DARK_GRAY, TFONT),
    ("", 6, False, BLACK, TFONT),
    ("Step 3: 应用 Bootstrapping 定理 → FHE", 15, True, GREEN, TFONT),
    ("  将噪声密文用新公钥加密 → 同态评估解密 → 得到新鲜密文", 13, False, DARK_GRAY, TFONT),
    ("  循环此过程可评估任意深度电路", 13, False, DARK_GRAY, TFONT),
    ("", 8, False, BLACK, TFONT),
    ("关键瓶颈：解密电路复杂度 >> SWHE的同态能力 → 需要特殊技巧！", 14, True, RED, TFONT),
])

# SLIDE 5: Section - "Motivation" -> "Re-linearization"
slide = prs.slides[4]
section_title(slide, "Re-linearization", "Part 2 - 核心创新I：重线性化")

# SLIDE 6: Content - Prior Work Limitations
slide = prs.slides[5]
content_title(slide, "前人工作的局限与本文化的突破")
clear_content_area(slide)
add_multiline_tb(slide, 2.8, 1.6, 9.0, 5.0, [
    ("局限 1：依赖理想格假设", 16, True, RED, TFONT),
    ("  Gentry SWHE需要理想格 → 理想格是特殊品种，认知不足", 13, False, DARK_GRAY, TFONT),
    ("  一般格的研究更深入(LLL, Ajtai, Micciancio...)但只支持加法", 12, False, MED_GRAY, TFONT),
    ("", 6, False, BLACK, TFONT),
    ("局限 2：Squashing + 稀疏子集和假设", 16, True, RED, TFONT),
    ("  解密电路太复杂 → 需Squashing人为降低 → 引入强假设", 13, False, DARK_GRAY, TFONT),
    ("  这个额外假设是Gentry方案及所有后续方案的主要缺陷", 12, False, MED_GRAY, TFONT),
    ("", 6, False, BLACK, TFONT),
    ("本文突破：", 16, True, GREEN, TFONT),
    ("  1. Re-linearization → SWHE from LWE (去掉理想格)", 14, True, MED_BLUE, TFONT),
    ("  2. Dimension-Modulus Reduction → Bootstrapping (去掉Squashing)", 14, True, MED_BLUE, TFONT),
])

# SLIDE 7: Section - "Overview" -> keep as is, or change
# Let's repurpose "Overview" as a content slide for LWE basics
slide = prs.slides[6]
section_title(slide, "LWE Encryption\nFoundation", "LWE加密基础：Regev方案 [STOC'05]")

# SLIDE 8: Content - flow diagram, repurpose for LWE + multiplication problem
slide = prs.slides[7]
content_title(slide, "LWE加密 → 线性函数视角 → 同态乘法难题")
clear_content_area(slide)
add_multiline_tb(slide, 2.8, 1.4, 9.0, 5.5, [
    ("Regev 加密 (单比特):", 16, True, DARK_BLUE, TFONT),
    ("  Secret Key: s in Z_q^n", 13, False, BLACK, MFONT),
    ("  Enc(mu): c = (a, b = <a,s> + 2e + mu)", 13, True, BLACK, MFONT),
    ("  Dec(c): (b - <a,s> mod q) mod 2 = mu", 13, False, BLACK, MFONT),
    ("", 6, False, BLACK, TFONT),
    ("解密视角转换 → 密文 = 关于密钥的线性函数:", 15, True, MED_BLUE, TFONT),
    ("  f_{a,b}(x) = b - <a,x> (mod q),   Dec = f_{a,b}(s) mod 2", 13, False, BLACK, MFONT),
    ("", 6, False, BLACK, TFONT),
    ("同态加法 ✓ 自然成立:", 15, True, GREEN, TFONT),
    ("  f1 + f2 = f_{a1+a2, b1+b2}  → 仍是线性!", 13, False, BLACK, MFONT),
    ("", 6, False, BLACK, TFONT),
    ("同态乘法 ✗ 密文膨胀:", 15, True, RED, TFONT),
    ("  f1 * f2 = phi(x) = h0 + Sum hi*x[i] + Sum hij*x[i]x[j]  (二次!)", 13, True, RED, MFONT),
    ("  密文大小: n+1 → O(n^2) → 一次乘法就不可用!", 13, False, RED, TFONT),
])

# SLIDE 9: Section - "Method" -> "Dimension-Modulus Reduction"
slide = prs.slides[8]
section_title(slide, "Dimension-Modulus\nReduction", "Part 3 - 核心创新II：维度-模缩减")

# SLIDE 10: Content - Re-linearization details
slide = prs.slides[9]
content_title(slide, "Re-linearization 技术详解")
clear_content_area(slide)
add_multiline_tb(slide, 2.8, 1.4, 9.0, 5.5, [
    ("核心思想：发布二次项 s[i]s[j] 在旧密钥下的伪加密 → 在新密钥下转化为线性组合", 14, True, MED_BLUE, TFONT),
    ("", 6, False, BLACK, TFONT),
    ("Step 1 - 密钥链：生成 s0, s1, ..., sL，每层支持一次乘法", 14, True, DARK_BLUE, TFONT),
    ("Step 2 - 发布伪加密参数 Psi：", 14, True, DARK_BLUE, TFONT),
    ("  b_{ell,i,j,tau} = <a, s_ell> + 2e + 2^tau * s_{ell-1}[i] * s_{ell-1}[j]", 12, False, BLACK, MFONT),
    ("Step 3 - 重线性化：Sum hij,tau * (b - <a, s_ell>) = 线性函数 in s_ell !", 14, True, DARK_BLUE, TFONT),
    ("", 6, False, BLACK, TFONT),
    ("效果：将二次密文 O(n^2) 压回线性 O(n)，每层乘法只需 ~n^2 log q 个辅助参数", 14, True, GREEN, TFONT),
    ("代价：hij 必须比特分解 (hij,tau in {0,1}) 以控制噪声增长", 13, False, DARK_GRAY, TFONT),
])

# SLIDE 11: Content - More re-linearization
slide = prs.slides[10]
content_title(slide, "Re-linearization: 成就与局限")
clear_content_area(slide)
add_multiline_tb(slide, 2.6, 1.4, 4.5, 5.0, [
    ("[成就] 取得的突破", 18, True, GREEN, TFONT),
    ("", 6, False, BLACK, TFONT),
    ("1. 从标准LWE构造出SWHE", 14, True, DARK_BLUE, TFONT),
    ("   摆脱理想格依赖!", 12, False, MED_GRAY, TFONT),
    ("2. 支持多项式深度运算", 14, True, DARK_BLUE, TFONT),
    ("   深度 L ~ epsilon*log n", 12, False, MED_GRAY, MFONT),
    ("   密文大小保持 O(n)", 12, False, MED_GRAY, TFONT),
    ("3. 密钥切换的雏形", 14, True, DARK_BLUE, TFONT),
    ("   BGV/GSW/CKKS的基础", 12, False, MED_GRAY, TFONT),
])
add_multiline_tb(slide, 7.4, 1.4, 4.5, 5.0, [
    ("[局限] 仍是 Somewhat HE", 18, True, RED, TFONT),
    ("", 6, False, BLACK, TFONT),
    ("解密电路: degree > max(n, log q)", 14, True, RED, MFONT),
    ("SWHE能力: degree ~ n^epsilon", 14, False, DARK_GRAY, MFONT),
    ("n^epsilon << max(n, log q)", 14, True, RED, TFONT),
    ("→ 无法满足自举条件!", 14, True, RED, TFONT),
    ("", 8, False, BLACK, TFONT),
    ("需要新工具降低解密复杂度", 14, True, MED_BLUE, TFONT),
    ("→ Dimension-Modulus Reduction!", 14, True, MED_BLUE, TFONT),
])

# SLIDE 12: (was empty) The Bootstrapping Bottleneck
slide = prs.slides[11]
content_title(slide, "Bootstrapping 瓶颈与维度-模缩减思路")
clear_content_area(slide)
add_multiline_tb(slide, 2.6, 1.4, 9.0, 5.5, [
    ("瓶颈：解密电路 degree >= max(n, log q)  >>  n^epsilon (SWHE能力)", 15, True, RED, TFONT),
    ("", 6, False, BLACK, TFONT),
    ("路径 A (Gentry'09): Squashing → 人为降低解密degree → 引入稀疏子集和假设 ✗", 13, False, DARK_GRAY, TFONT),
    ("路径 B (本文): Dimension-Modulus Reduction  → 缩减密文参数 → 无需额外假设 ✓", 14, True, GREEN, TFONT),
    ("", 6, False, BLACK, TFONT),
    ("维度缩减 (n→k): 选择短密钥 s_hat in Z_p^k (k<<n), Re-linearization照样工作", 14, True, MED_BLUE, TFONT),
    ("模缩减 (q→p): scale(x) = round((p/q)*x), 舍入误差<=1/2 在噪声预算内可控", 14, True, MED_BLUE, TFONT),
    ("", 6, False, BLACK, TFONT),
    ("效果：(n, log q) → (k, log p) → max(k, log p) 很小 → 解密可在SWHE能力内评估!", 15, True, GREEN, TFONT),
])

# SLIDE 13: Section - "Evaluation" -> repurpose as BTS scheme section
# Actually let's keep evaluations section for BTS
slide = prs.slides[12]
section_title(slide, "BTS Scheme\nConstruction", "BTS方案完整构造")

# SLIDE 14: Content - BTS details
slide = prs.slides[13]
content_title(slide, "BTS 方案: KeyGen / Enc / Dec / Eval")
clear_content_area(slide)
add_multiline_tb(slide, 2.8, 1.4, 9.0, 5.5, [
    ("KeyGen: 长密钥链s0..sL+短密钥s_hat+Psi(重线性化)+Psi_hat(桥接)+pk=(A,b=As0+2e)", 12, False, BLACK, MFONT),
    ("Enc: r<-{0,1}^m, v=A^T*r, w=b^T*r+mu, 输出 ((v,w), tag=0)", 12, False, BLACK, MFONT),
    ("Dec: mu* = (w_hat - <v_hat,s_hat> mod p) mod 2", 12, False, BLACK, MFONT),
    ("", 6, False, BLACK, TFONT),
    ("同态评估不变量: w - <v,s_ell> = mu + 2e (mod q), |e| < q/4 保证正确解密", 14, True, MED_BLUE, TFONT),
    ("", 6, False, BLACK, TFONT),
    ("加法Gate: c_add = ((Sum vi, Sum wi), ell) → 噪声线性增长", 14, True, DARK_BLUE, TFONT),
    ("乘法Gate (3步):", 14, True, DARK_BLUE, TFONT),
    ("  1. 构造二次型 Phi(x) = (w-<v,x>)(w'-<v',x>) = Sum hij * x[i]x[j]", 12, False, BLACK, MFONT),
    ("  2. 比特分解 hij = Sum_tau hij,tau * 2^tau  (hij,tau in {0,1})", 12, False, BLACK, MFONT),
    ("  3. 重线性化: v_mult=Sum hij,tau*a; w_mult=Sum hij,tau*b → tag=ell+1", 12, False, BLACK, MFONT),
    ("最后一步-维度模缩减: 使用Psi_hat将(n,q)密文转为(k,p)短密文", 14, True, MED_BLUE, TFONT),
])

# SLIDE 15: Section - "Conclusion" -> "PIR Application"
slide = prs.slides[14]
section_title(slide, "Application:\nNear-Optimal PIR", "Part 4 - 应用：近最优私有信息检索")

# SLIDE 16: Content - Contribution -> change to PIR + Bootstrapping
slide = prs.slides[15]
content_title(slide, "Bootstrapping达成 + PIR应用")
clear_content_area(slide)
add_multiline_tb(slide, 2.6, 1.4, 4.5, 5.0, [
    ("Bootstrapping达成", 18, True, GREEN, TFONT),
    ("", 6, False, BLACK, TFONT),
    ("参数: k=kappa, n=k^4, q=2^sqrt(n)", 13, False, BLACK, MFONT),
    ("p=poly(k), L=(1/3)log n", 13, False, BLACK, MFONT),
    ("解密电路: O(log k+log log p)深度", 13, False, DARK_GRAY, TFONT),
    ("BTS能力: ~n^epsilon = k^{c*epsilon}", 13, False, DARK_GRAY, TFONT),
    ("选c够大 → 可自举!", 14, True, GREEN, TFONT),
    ("安全性: DLWE_{n,q} + DLWE_{k,p}", 13, False, DARK_GRAY, TFONT),
])
add_multiline_tb(slide, 7.4, 1.4, 4.5, 5.0, [
    ("近最优PIR", 18, True, MED_BLUE, TFONT),
    ("", 6, False, BLACK, TFONT),
    ("对称加密索引 → O(logN)查询", 13, False, DARK_GRAY, TFONT),
    ("BTS短密文: O(k log k)比特", 13, False, DARK_GRAY, TFONT),
    ("响应: O(logN*polyloglogN)", 14, True, GREEN, TFONT),
    ("", 6, False, BLACK, TFONT),
    ("首个基于LWE的polylog PIR!", 14, True, MED_BLUE, TFONT),
    ("接近理论下界Omega(logN)", 13, False, MED_GRAY, TFONT),
])

# SLIDE 17: Content - empty, use as Summary
slide = prs.slides[16]
content_title(slide, "总结与贡献")
clear_content_area(slide)
add_multiline_tb(slide, 2.6, 1.4, 9.0, 5.5, [
    ("贡献1: Re-linearization", 18, True, MED_BLUE, TFONT),
    ("  → SWHE from 标准LWE (去掉理想格) → 密钥切换范式 → BGV/GSW/CKKS等的基础", 13, False, DARK_GRAY, TFONT),
    ("贡献2: Dimension-Modulus Reduction", 18, True, MED_BLUE, TFONT),
    ("  → 替代Squashing (去掉稀疏子集和) → 纯LWE Bootstrapping → 短密文 → 近最优PIR", 13, False, DARK_GRAY, TFONT),
    ("", 6, False, BLACK, TFONT),
    ("开放问题：循环安全性(去最后额外假设) / 常数摊销PIR / 实际效率优化", 14, True, DARK_BLUE, TFONT),
    ("后续影响：BGV(层次型FHE), GSW(矩阵型FHE), CKKS(浮点FHE), TFHE(快速自举) 等均基于本文技术", 13, False, MED_GRAY, TFONT),
])

# ================================================================
# SAVE
# ================================================================
prs.save(OUTPUT)
print(f"DONE! Saved to: {OUTPUT}")
print(f"Total slides: {len(prs.slides)}")
