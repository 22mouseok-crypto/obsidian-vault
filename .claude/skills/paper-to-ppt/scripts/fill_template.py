#!/usr/bin/env python3
"""
PPT 模板处理脚本 — 风格提取、布局分析、内容填充。

用法:
    # 提取模板风格
    python fill_template.py template.pptx --extract-style template_style.json

    # 分析可用布局
    python fill_template.py template.pptx --analyze-layouts

    # 根据 slide_plan.json 填充内容
    python fill_template.py template.pptx --fill slide_plan.json \
        --paper-outline paper_outline.json \
        --figures-dir ./figures/ \
        --output 汇报.pptx
"""

import argparse
import json
import os
import sys
from copy import deepcopy
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    print("Error: python-pptx is required. Install: pip install python-pptx")
    sys.exit(1)

try:
    from PIL import Image
except ImportError:
    Image = None


# ============================================================
# 模式一：提取模板风格
# ============================================================

def extract_style(prs: Presentation, output_path: str):
    """提取 PPT 模板的 Slide Master 设计系统。"""
    style_info = {
        "slide_width": prs.slide_width,
        "slide_height": prs.slide_height,
        "slide_masters": [],
        "layouts": [],
    }

    for mi, master in enumerate(prs.slide_masters):
        master_info = {"index": mi, "name": master.name if hasattr(master, 'name') else f"Master-{mi}"}

        # 提取主题配色
        try:
            theme = master.slide_layouts[0] if master.slide_layouts else None
        except Exception:
            theme = None
        master_info["layouts"] = []

    # 遍历所有布局
    for li, layout in enumerate(prs.slide_layouts):
        layout_info = {
            "index": li,
            "name": layout.name,
            "placeholders": [],
        }

        for ph in layout.placeholders:
            phf = ph.placeholder_format
            ph_info = {
                "idx": phf.idx,
                "type": str(phf.type),
                "name": ph.name,
                "width": ph.width,
                "height": ph.height,
                "left": ph.left,
                "top": ph.top,
            }
            # 提取占位符中的默认字体信息
            if ph.has_text_frame:
                for para in ph.text_frame.paragraphs:
                    for run in para.runs:
                        font = run.font
                        ph_info["font_name"] = font.name
                        ph_info["font_size"] = font.size / 12700 if font.size else None  # EMU → pt
                        try:
                            if font.color and font.color.rgb:
                                ph_info["font_color"] = str(font.color.rgb)
                        except Exception:
                            pass
                        break
                    break

            layout_info["placeholders"].append(ph_info)

        style_info["layouts"].append(layout_info)

    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(style_info, f, ensure_ascii=False, indent=2)

    print(f"✅ 模板风格已导出: {output_path}")
    print(f"   • 幻灯片尺寸: {style_info['slide_width']}×{style_info['slide_height']} EMU")
    print(f"   • 可用布局数: {len(style_info['layouts'])}")
    for layout in style_info["layouts"]:
        print(f"     [{layout['index']}] {layout['name']} ({len(layout['placeholders'])} placeholders)")


# ============================================================
# 模式二：分析布局
# ============================================================

def analyze_layouts(prs: Presentation):
    """详细分析每个布局的结构。"""
    print("=" * 60)
    print("PPT 模板布局分析")
    print("=" * 60)
    print(f"幻灯片尺寸: {prs.slide_width}×{prs.slide_height} EMU")
    print(f"  → 约 {prs.slide_width/914400:.1f}\"×{prs.slide_height/914400:.1f}\"")
    print()

    for li, layout in enumerate(prs.slide_layouts):
        print(f"[布局 {li}] {layout.name}")
        print(f"  占位符数量: {len(layout.placeholders)}")

        for ph in layout.placeholders:
            phf = ph.placeholder_format
            has_text = "📝" if ph.has_text_frame else ""
            has_img = "🖼️" if hasattr(ph, 'image') else ""
            print(f"    idx={phf.idx:>3} | {phf.type} | '{ph.name}' {has_text}{has_img}")
            print(f"           位置: left={ph.left}, top={ph.top}, {ph.width}×{ph.height}")

            if ph.has_text_frame:
                preview = ph.text_frame.text[:80].replace("\n", " ")
                if preview:
                    print(f"           预置文本: \"{preview}\"")
        print()


# ============================================================
# 模式三：填充幻灯片
# ============================================================

def fill_slides(
    template_path: str,
    slide_plan_path: str,
    paper_outline_path: str,
    figures_dir: str,
    output_path: str,
):
    """根据 slide_plan.json 逐页填充 PPT 模板。"""
    prs = Presentation(template_path)

    # 加载数据
    with open(slide_plan_path, "r", encoding="utf-8") as f:
        plan = json.load(f)

    paper_outline = None
    if paper_outline_path and os.path.exists(paper_outline_path):
        with open(paper_outline_path, "r", encoding="utf-8") as f:
            paper_outline = json.load(f)

    # 加载图表文件列表
    figure_files = {}
    if figures_dir and os.path.isdir(figures_dir):
        for fname in os.listdir(figures_dir):
            if fname.lower().endswith((".png", ".jpg", ".jpeg", ".gif", ".bmp")):
                figure_files[fname] = os.path.join(figures_dir, fname)

    # 清空模板中已有的幻灯片（保留 Slide Master）
    _remove_all_slides(prs)

    # 逐页创建
    slides_plan = plan.get("slides", [])
    total = len(slides_plan)

    for i, slide_plan in enumerate(slides_plan):
        page_num = slide_plan.get("page", i + 1)
        layout_idx = slide_plan.get("layout_index", 1)
        content = slide_plan.get("content", {})

        print(f"📝 处理第 {page_num}/{total} 页: {content.get('title', 'Untitled')[:60]}")

        # 选择布局
        if layout_idx >= len(prs.slide_layouts):
            print(f"   ⚠️ 布局索引 {layout_idx} 超出范围，使用默认布局 0")
            layout_idx = 0

        layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(layout)

        # 填充内容
        _fill_slide_content(slide, content, figure_files, paper_outline)

    # 保存
    prs.save(output_path)
    print(f"\n✅ PPT 生成完成: {output_path}")
    print(f"   • 总页数: {total}")
    print(f"   • 保留模板: {Path(template_path).name}")


def _remove_all_slides(prs: Presentation):
    """删除所有幻灯片但保留 Slide Master。"""
    # 通过 XML 操作删除所有 slide 引用
    sldIdLst = prs.presentation.sldIdLst
    for sldId in list(sldIdLst):
        sldIdLst.remove(sldId)


def _fill_slide_content(slide, content: dict, figure_files: dict, paper_outline: dict):
    """填充单页幻灯片内容。"""
    title = content.get("title", "")
    subtitle = content.get("subtitle", "")
    bullets = content.get("bullets", [])
    extra = content.get("extra", "")
    notes = content.get("notes", "")
    image_ref = content.get("image", "")
    layout_type = content.get("_type", "content")

    # 1. 填充标题
    if title and slide.shapes.title:
        _safe_set_text(slide.shapes.title, title)

    # 2. 填充副标题（如果有 subtitle placeholder）
    if subtitle:
        for shape in slide.placeholders:
            try:
                if shape.placeholder_format.idx == 1 and shape.has_text_frame:
                    if "subtitle" in shape.name.lower() or layout_type == "title":
                        _safe_set_text(shape, subtitle)
                        break
            except Exception:
                continue

    # 3. 填充要点的 body placeholder
    if bullets:
        body_shape = None
        # 尝试找到 content/body placeholder
        for shape in slide.placeholders:
            try:
                phf = shape.placeholder_format
                if phf.idx == 1 and shape.has_text_frame:
                    name_lower = shape.name.lower()
                    if "body" in name_lower or "content" in name_lower or "text" in name_lower:
                        body_shape = shape
                        break
            except Exception:
                continue

        # 如果没找到 body placeholder，找任意可用的文本 placeholder
        if body_shape is None:
            for shape in slide.placeholders:
                if shape.has_text_frame and shape != slide.shapes.title:
                    body_shape = shape
                    break

        if body_shape:
            tf = body_shape.text_frame
            tf.clear()
            for i, bullet in enumerate(bullets):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = bullet
                p.level = 0
                p.space_after = Pt(6)
        else:
            # 动态创建文本框
            slide_width = slide.part.slide_layout.slide_master.slide_width if hasattr(slide, 'part') else Emu(9144000)
            left = Inches(0.8)
            top = Inches(1.8)
            width = slide_width - Inches(1.6)
            height = Inches(4.5)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True
            for i, bullet in enumerate(bullets):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                run = p.add_run()
                run.text = f"• {bullet}"
                run.font.size = Pt(14)
                p.space_after = Pt(6)

    # 4. 嵌入图片
    if image_ref and image_ref in figure_files:
        img_path = figure_files[image_ref]
        try:
            # 放在右下区域
            left = Inches(5.5)
            top = Inches(3.5)
            width = Inches(3.5)
            if Image:
                with Image.open(img_path) as img:
                    aspect = img.width / img.height
                    height = int(width / aspect)
            else:
                height = Inches(2.5)
            slide.shapes.add_picture(img_path, left, top, width, height)
            print(f"   🖼️  嵌入图片: {image_ref}")
        except Exception as e:
            print(f"   ⚠️ 图片嵌入失败: {e}")

    # 5. 添加备注
    if notes:
        try:
            slide.notes_slide.notes_text_frame.text = notes
        except Exception:
            pass


def _safe_set_text(shape, text: str):
    """安全设置文本，保留首个 run 的格式。"""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    # 保留原有格式：替换第一个 run 的文本
    if tf.paragraphs and tf.paragraphs[0].runs:
        tf.paragraphs[0].runs[0].text = text
    else:
        tf.paragraphs[0].text = text


# ============================================================
# 入口
# ============================================================

def main():
    parser = argparse.ArgumentParser(description="PPT 模板处理工具")
    parser.add_argument("template", help="PPT 模板文件路径 (.pptx)")
    parser.add_argument("--extract-style", help="导出模板风格到 JSON 文件")
    parser.add_argument("--analyze-layouts", action="store_true", help="分析模板布局")
    parser.add_argument("--fill", help="slide_plan.json 文件路径")
    parser.add_argument("--paper-outline", help="论文大纲 JSON 文件（可选，用于引用论文数据）")
    parser.add_argument("--figures-dir", help="论文图表导出目录（可选）")
    parser.add_argument("--output", "-o", default="output.pptx", help="输出 PPT 文件路径")
    args = parser.parse_args()

    if not os.path.exists(args.template):
        print(f"Error: Template not found: {args.template}")
        sys.exit(1)

    prs = Presentation(args.template)

    if args.extract_style:
        extract_style(prs, args.extract_style)

    if args.analyze_layouts:
        analyze_layouts(prs)

    if args.fill:
        fill_slides(
            args.template,
            args.fill,
            args.paper_outline,
            args.figures_dir,
            args.output,
        )

    # 如果什么都不做，默认运行分析
    if not args.extract_style and not args.analyze_layouts and not args.fill:
        analyze_layouts(prs)


if __name__ == "__main__":
    main()
